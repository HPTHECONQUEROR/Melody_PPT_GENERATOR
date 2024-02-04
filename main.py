import io
import re
import random
import threading
import requests
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
import os
from pptx.dml.color import RGBColor
from concurrent.futures import ThreadPoolExecutor, as_completed

command = "create a ppt on the topic terminator movie  with 3 slides"

# Extracting the number of slides using re
slides_match = re.search(r'\d+', command)
slides = int(slides_match.group()) if slides_match else None
if slides<=5:
    pass
else:
    slides =5
    print("sorry the slides limit is 5")
headers = {"Authorization": "Bearer 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"}

# Extracting the topic from the command
topic_start = command.find("topic") + len("topic")
topic_end = command.find("with")
topic = command[topic_start:topic_end].strip() if topic_start < topic_end else None

filename = f"{topic}.pptx"
count = 1
while os.path.exists(filename):
    filename = f"{topic}_{count}.pptx"
    count += 1

## makes heading and body

API_URL_TEXT = "https://api-inference.huggingface.co/models/NousResearch/Nous-Hermes-2-Mixtral-8x7B-DPO"
headers = {"Authorization": "Bearer <REPLACE YOUR OWN KEY>"}
command_ip = f" {slides},{topic} "
with open("textfiles/topics.txt", "r") as file:
    prompt_ = file.read()
    file.close()
input_text_prompt = f"\n<user> {command_ip} <u/> <assistant> "
prompt = prompt_ + input_text_prompt
for i in range(slides):
    response = requests.post(API_URL_TEXT, headers=headers, json={"inputs": prompt, })
    output = str(response.json()[0]["generated_text"])

topics = output[output.rfind(input_text_prompt) + len(input_text_prompt) - 1:]
topics = topics[:topics.find("<a/>")]
topics = topics.split(",")
with open("textfiles/prompt.txt") as file:
    topic_prompt = file.read()
    file.close()
topic_list = []


def send_api_request(topic):
    input_topic_prompt = f"\n<user> {topic} <u/> <assistant> "
    response = requests.post(API_URL_TEXT, headers=headers, json={"inputs": input_topic_prompt})
    output_topic = str(response.json()[0]["generated_text"])
    output_topic = output_topic[:output_topic.rfind("<a/>")]
    append_topic = output_topic[output_topic.rfind(input_topic_prompt) + len(input_topic_prompt):]
    return append_topic


with ThreadPoolExecutor(max_workers=len(topics)) as executor:
    # Submit tasks and retrieve futures
    futures = {executor.submit(send_api_request, topic): topic for topic in topics}

    topic_list = []
    for future in as_completed(futures):
        topic = futures[future]
        try:
            result = future.result()
            topic_list.append(result)
        except Exception as e:
            print(f"Error occurred for topic {topic}: {e}")
topics = topics[:slides]
topic_list = topic_list[:slides]

##ends here


with open("textfiles/tags.txt", "r") as file:
    tags_prompt = file.read()
    file.close()
input_text = f"\n<user> {slides},{topic}<u/><assistant>"
tags_prompt += input_text

API_URL = "https://api-inference.huggingface.co/models/NousResearch/Nous-Hermes-2-Mixtral-8x7B-DPO"
response = requests.post(API_URL, headers=headers, json={"inputs": tags_prompt, })
output = str(response.json()[0]["generated_text"])
# print(output)

tags = output[output.rfind(f"{input_text}") + len(input_text):]
tag_list = re.split(r'(\d+\.)', tags)[1:]

# Combining the strings into pairs
tag_pairs = [f"{tag_list[i]}{tag_list[i + 1]}" for i in range(0, len(tag_list), 2)]
tag_pairs = tag_pairs[:slides]

generated_images = []
API_URL = "https://api-inference.huggingface.co/models/goofyai/cyborg_style_xl"
headers = {"Authorization": "Bearer <REPLACE YOUR OWN KEY>"}


def query(payload):
    response = requests.post(API_URL, headers=headers, json=payload)
    return response.content


generated_images = []


def api_call(tag, index):
    image_bytes = query({"inputs": tag})
    image = Image.open(io.BytesIO(image_bytes))
    filename = f"{topic}{index + 1}.png"
    image.save(filename)
    generated_images.append(filename)


threads = []
for i in range(min(slides, len(tags_prompt))):
    thread = threading.Thread(target=api_call, args=(tag_pairs[i], i))
    threads.append(thread)
    thread.start()

for thread in threads:
    thread.join()


presentation = Presentation()

presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

for i in range(slides):
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)

    background = slide.background
    fill = background.fill
    fill.solid()
    rgb_combinations = [
        RGBColor(0, 0, 0),
        RGBColor(55, 65, 74),
        RGBColor(89, 46, 83),
        RGBColor(1, 23, 49)

    ]


    random_rgb = random.choice(rgb_combinations)


    fill.fore_color.rgb = random_rgb

    # Calculate the size and position for the image
    image_path = generated_images[i]
    image = Image.open(image_path)
    image_width, image_height = image.size
    target_width = Inches(9)
    target_height = (target_width / image_width) * image_height
    left_image = Inches(-2.5) if i % 2 == 0 else Inches(9.5)
    top_image = Inches(0)

    pic = slide.shapes.add_picture(image_path, left_image, top_image, width=target_width, height=target_height)

    # Add heading
    left_heading = Inches(5.5) if i % 2 == 0 else Inches(0)
    top_heading = Inches(0)
    heading_shape = slide.shapes.add_textbox(left_heading, top_heading, Inches(10), Inches(2))
    heading_frame = heading_shape.text_frame
    heading = heading_frame.add_paragraph()
    heading.text = topics[i]
    heading.font.color.rgb = RGBColor(255, 255, 255)
    heading.font.size = Pt(38)
    heading.font.bold = True
    heading.font.name = 'Times New Roman'
    heading.word_wrap = True

    left_body = Inches(5.9) if i % 2 == 0 else Inches(0)
    top_body = Inches(2)
    width_body = Inches(8)
    height_body = Inches(3)

    body_shape = slide.shapes.add_textbox(left_body, top_body, Inches(10), Inches(7))
    body_shape.text_frame.word_wrap = True
    body_frame = body_shape.text_frame
    body = body_frame.add_paragraph()
    body.text = topic_list[i]
    body.font.size = Pt(29)
    body.font.name = 'Times New Roman'
    body.font.color.rgb = RGBColor(255, 255, 255)

    body_shape.line.width = Pt(5)
    body_shape.line.color.rgb = RGBColor(
        random.randint(200, 255),
        random.randint(200, 255),
        random.randint(200, 255)
    )

    if i == 0:
        small_image_path = "img.png"
        small_image_width = Inches(2.58)
        small_image_height = Inches(0.52)
        left_small_image = Inches(13)
        top_small_image = Inches(8)
        slide.shapes.add_picture(small_image_path, left_small_image, top_small_image, width=small_image_width,
                                 height=small_image_height)





presentation.save(f"ppts/{filename}")